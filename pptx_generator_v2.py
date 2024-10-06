import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    lv = len(hex_color)
    return RGBColor(
        *(int(hex_color[i : i + lv // 3], 16) for i in range(0, lv, lv // 3))
    )


def create_pptx_from_json(data, uuid=None):
    prs = Presentation()

    # Get logo image
    logo_img = None
    if data.get("logoURL"):
        try:
            response = requests.get(data["logoURL"])
            logo_img = BytesIO(response.content)
        except:
            pass  # Handle error or leave logo_img as None

    # Get background color
    bg_color = None
    if data.get("color"):
        try:
            bg_color = hex_to_rgb(data["color"])
        except:
            pass

    # Iterate over list items
    for item in data.get("list", []):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background color if available
        if bg_color:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color

        # Add logo if available
        if logo_img:
            try:
                slide.shapes.add_picture(
                    logo_img, Inches(0.5), Inches(0.5), width=Inches(1)
                )
            except:
                pass

        if item["type"] == "HERO":
            # Add image if imageURL is provided
            if item.get("imageURL"):
                try:
                    response = requests.get(item["imageURL"])
                    image = BytesIO(response.content)
                    slide.shapes.add_picture(
                        image, Inches(1), Inches(1), width=Inches(8), height=Inches(4)
                    )
                except:
                    pass  # Handle error or skip image

            # Add title
            if item.get("title"):
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(5.5), Inches(8), Inches(1)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = item["title"]
                font = run.font
                font.size = Pt(32)
                font.bold = True
                p.alignment = PP_ALIGN.CENTER

            # Add subtitle
            if item.get("subtitle"):
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(6.5), Inches(8), Inches(1)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = item["subtitle"]
                font = run.font
                font.size = Pt(24)
                p.alignment = PP_ALIGN.CENTER

        elif item["type"] == "FEATURES":
            # Add title
            if item.get("title"):
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(1), Inches(8), Inches(1)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = item["title"]
                font = run.font
                font.size = Pt(32)
                font.bold = True
                p.alignment = PP_ALIGN.CENTER

            # Arrange features in columns
            features = item.get("features", [])
            num_features = len(features)
            if num_features > 0:
                column_width = 8 / num_features
                for i, feature in enumerate(features):
                    x_position = 1 + i * column_width
                    txBox = slide.shapes.add_textbox(
                        Inches(x_position), Inches(2), Inches(column_width), Inches(5)
                    )
                    tf = txBox.text_frame

                    # Add emoji and title
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = f"{feature.get('emoji', '')} {feature.get('title', '')}"
                    font = run.font
                    font.size = Pt(24)
                    font.bold = True
                    p.alignment = PP_ALIGN.CENTER

                    # Add description
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = feature.get("description", "")
                    font = run.font
                    font.size = Pt(18)
                    p.alignment = PP_ALIGN.CENTER

        elif item["type"] == "BENEFITS":
            # Add title
            if item.get("title"):
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(1), Inches(8), Inches(1)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = item["title"]
                font = run.font
                font.size = Pt(32)
                font.bold = True
                p.alignment = PP_ALIGN.CENTER

            # Arrange benefits in columns
            benefits = item.get("benefits", [])
            num_benefits = len(benefits)
            if num_benefits > 0:
                column_width = 8 / num_benefits
                for i, benefit in enumerate(benefits):
                    x_position = 1 + i * column_width
                    txBox = slide.shapes.add_textbox(
                        Inches(x_position), Inches(2), Inches(column_width), Inches(5)
                    )
                    tf = txBox.text_frame

                    # Add emoji and name
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = f"{benefit.get('emoji', '')} {benefit.get('name', '')}"
                    font = run.font
                    font.size = Pt(24)
                    font.bold = True
                    p.alignment = PP_ALIGN.CENTER

                    # Add description
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = benefit.get("description", "")
                    font = run.font
                    font.size = Pt(18)
                    p.alignment = PP_ALIGN.CENTER

        elif item["type"] == "EXPLANATION":
            explanations = item.get("explanations", [])
            num_explanations = len(explanations)
            y_position = 1.5
            for explanation in explanations:
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(y_position), Inches(8), Inches(2)
                )
                tf = txBox.text_frame

                # Add emoji and title
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = (
                    f"{explanation.get('emoji', '')} {explanation.get('title', '')}"
                )
                font = run.font
                font.size = Pt(24)
                font.bold = True

                # Add description
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = explanation.get("description", "")
                font = run.font
                font.size = Pt(18)

                y_position += 2.5  # Adjust y position

        elif item["type"] == "TESTIMONIALS":
            testimonials = item.get("testimonials", [])
            num_testimonials = len(testimonials)
            testimonials_per_slide = 2
            for i in range(0, num_testimonials, testimonials_per_slide):
                if i > 0:
                    # For additional slides
                    slide = prs.slides.add_slide(slide_layout)
                    # Set background color if available
                    if bg_color:
                        fill = slide.background.fill
                        fill.solid()
                        fill.fore_color.rgb = bg_color
                    # Add logo if available
                    if logo_img:
                        slide.shapes.add_picture(
                            logo_img, Inches(0.5), Inches(0.5), width=Inches(1)
                        )

                y_position = 1.5
                for j in range(testimonials_per_slide):
                    if i + j >= num_testimonials:
                        break
                    testimonial = testimonials[i + j]
                    txBox = slide.shapes.add_textbox(
                        Inches(1), Inches(y_position), Inches(8), Inches(2)
                    )
                    tf = txBox.text_frame

                    # Add testimonial text
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = f'"{testimonial.get("testimonial", "")}"'
                    font = run.font
                    font.size = Pt(18)
                    font.italic = True

                    # Add name
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = f"- {testimonial.get('firstName', '')} {testimonial.get('lastName', '')}"
                    font = run.font
                    font.size = Pt(16)
                    font.bold = True

                    y_position += 3  # Adjust y position

        elif item["type"] == "CTA":
            # Add headline
            if item.get("headline"):
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(2), Inches(8), Inches(1)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = item["headline"]
                font = run.font
                font.size = Pt(32)
                font.bold = True
                p.alignment = PP_ALIGN.CENTER

            # Add description
            if item.get("description"):
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(3.5), Inches(8), Inches(1)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = item["description"]
                font = run.font
                font.size = Pt(24)
                p.alignment = PP_ALIGN.CENTER

            # Add link and homepageLink
            y_position = 5
            if item.get("link"):
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(y_position), Inches(8), Inches(0.5)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f"Call to Action: {item['link']}"
                font = run.font
                font.size = Pt(18)
                font.color.rgb = RGBColor(0, 0, 255)  # Blue color for link
                p.alignment = PP_ALIGN.CENTER
                y_position += 0.75

            if item.get("homepageLink"):
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(y_position), Inches(8), Inches(0.5)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f"Homepage: {item['homepageLink']}"
                font = run.font
                font.size = Pt(18)
                font.color.rgb = RGBColor(0, 0, 255)  # Blue color for link
                p.alignment = PP_ALIGN.CENTER

    # Save the presentation to a file
    output_filename = f"decks/{uuid}.pptx"
    prs.save(output_filename)
    return output_filename
