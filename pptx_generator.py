from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from io import BytesIO


# Helper function to add title and subtitle
def add_title_slide(prs, title, subtitle, imageURL=None):
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Add an image if URL is provided
    if imageURL:
        response = requests.get(imageURL)
        image_stream = BytesIO(response.content)
        try:
            slide.shapes.add_picture(
                image_stream, Inches(1), Inches(1), width=Inches(3)
            )
        except:
            pass


# Helper function to add content slide with bullet points
def add_content_slide(prs, title, items, content_type):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title

    # Create bullet point text
    text_frame = body_shape.text_frame
    for item in items:
        p = text_frame.add_paragraph()
        p.text = f"{item.get('emoji', '')} {item.get('title', '') or item.get('name', '')}: {item.get('description', '')}"
        p.font.size = Pt(18)


# Helper function to add CTA slide
def add_cta_slide(prs, headline, description, link, homepageLink):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = headline

    # Add description and links
    text_frame = body_shape.text_frame
    text_frame.text = description
    p = text_frame.add_paragraph()
    p.text = f"More Info: {homepageLink}"
    p = text_frame.add_paragraph()
    p.text = f"Call to Action: {link}"


def create_pptx_from_json(data, uuid=None):
    # Initialize presentation
    prs = Presentation()

    # Process each item in the list
    for item in data["list"]:
        slide_type = item.get("type")

        if slide_type == "HERO":
            add_title_slide(
                prs,
                item.get("title", ""),
                item.get("subtitle", ""),
                item.get("imageURL"),
            )

        elif slide_type == "FEATURES":
            add_content_slide(
                prs, item.get("title", "Features"), item.get("features", []), slide_type
            )

        elif slide_type == "BENEFITS":
            add_content_slide(
                prs, item.get("title", "Benefits"), item.get("benefits", []), slide_type
            )

        elif slide_type == "EXPLANATION":
            add_content_slide(
                prs,
                item.get("title", "Explanation"),
                item.get("explanations", []),
                slide_type,
            )

        elif slide_type == "TESTIMONIALS":
            add_content_slide(
                prs, "Testimonials", item.get("testimonials", []), slide_type
            )

        elif slide_type == "CTA":
            add_cta_slide(
                prs,
                item.get("headline", ""),
                item.get("description", ""),
                item.get("link", ""),
                item.get("homepageLink", ""),
            )

    # Save the presentation to a file
    output_filename = f"{uuid}.pptx"
    prs.save(output_filename)
    return output_filename
