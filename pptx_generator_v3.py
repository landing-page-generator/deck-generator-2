import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from io import BytesIO


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    lv = len(hex_color)
    return RGBColor(
        *(int(hex_color[i : i + lv // 3], 16) for i in range(0, lv, lv // 3))
    )


def create_pptx_from_json(data, uuid=None):
    prs = Presentation()

    # Define a consistent theme color
    theme_color = None
    if data.get("color"):
        try:
            theme_color = hex_to_rgb(data["color"])
        except:
            pass
    else:
        # Default theme color
        theme_color = RGBColor(0, 112, 192)  # A shade of blue

    # Get logo image
    logo_img = None
    if data.get("logoURL"):
        try:
            response = requests.get(data["logoURL"])
            logo_img = BytesIO(response.content)
        except:
            pass  # Handle error or leave logo_img as None

    # Iterate over list items
    for item in data.get("list", []):
        if item["type"] == "HERO":
            slide_layout = prs.slide_layouts[0]  # Title Slide layout
            slide = prs.slides.add_slide(slide_layout)

            # Set background image if imageURL is provided
            if item.get("imageURL"):
                try:
                    response = requests.get(item["imageURL"])
                    image = BytesIO(response.content)
                    # Set the background image
                    fill = slide.background.fill
                    fill.solid()
                    image_part = prs.part.related_parts[prs.part.relate_to_image(image)]
                    fill.fore_color.type = MSO_THEME_COLOR.ACCENT_1
                    fill.fore_color._xFill.solidFill.blipFill = image_part.blob
                except:
                    pass  # Handle error or skip image
            else:
                # Set background color if image is not available
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = theme_color

            # Add title
            if item.get("title"):
                title_shape = slide.shapes.title
                title_shape.text = item["title"]
                title_shape.text_frame.paragraphs[0].font.size = Pt(44)
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(
                    255, 255, 255
                )  # White text
                title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Add subtitle
            if item.get("subtitle"):
                # Find the subtitle placeholder by type
                subtitle_placeholder = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                        subtitle_placeholder = shape
                        break
                if subtitle_placeholder:
                    subtitle_placeholder.text = item["subtitle"]
                    subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
                    subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = (
                        RGBColor(255, 255, 255)
                    )  # White text
                    subtitle_placeholder.text_frame.paragraphs[0].alignment = (
                        PP_ALIGN.CENTER
                    )
                else:
                    # If subtitle placeholder not found, create a textbox
                    left = Inches(1)
                    top = Inches(4)
                    width = prs.slide_width - Inches(2)
                    height = Inches(1)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = item["subtitle"]
                    tf.paragraphs[0].font.size = Pt(24)
                    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        elif item["type"] == "FEATURES":
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            # Add title
            if item.get("title"):
                title_shape = slide.shapes.title
                title_shape.text = item["title"]
                title_shape.text_frame.paragraphs[0].font.size = Pt(36)
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.color.rgb = theme_color

            # Access the content placeholder by type
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    body_shape = shape
                    break
            if body_shape is None:
                # Create a textbox if placeholder not found
                body_shape = slide.shapes.add_textbox(
                    Inches(1), Inches(2), Inches(8), Inches(5)
                )

            tf = body_shape.text_frame
            tf.clear()  # Remove any existing paragraphs
            for feature in item.get("features", []):
                p = tf.add_paragraph()
                p.text = f"{feature.get('emoji', '')} {feature.get('title', '')}: {feature.get('description', '')}"
                p.font.size = Pt(24)
                p.level = 0

        elif item["type"] == "BENEFITS":
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            # Add title
            if item.get("title"):
                title_shape = slide.shapes.title
                title_shape.text = item["title"]
                title_shape.text_frame.paragraphs[0].font.size = Pt(36)
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.color.rgb = theme_color

            # Access the content placeholder by type
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    body_shape = shape
                    break
            if body_shape is None:
                # Create a textbox if placeholder not found
                body_shape = slide.shapes.add_textbox(
                    Inches(1), Inches(2), Inches(8), Inches(5)
                )

            tf = body_shape.text_frame
            tf.clear()
            for benefit in item.get("benefits", []):
                p = tf.add_paragraph()
                p.text = f"{benefit.get('emoji', '')} {benefit.get('name', '')}: {benefit.get('description', '')}"
                p.font.size = Pt(24)
                p.level = 0

        elif item["type"] == "EXPLANATION":
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            # Title
            slide.shapes.title.text = "Explanation"
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = theme_color

            # Access the content placeholder
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    body_shape = shape
                    break
            if body_shape is None:
                body_shape = slide.shapes.add_textbox(
                    Inches(1), Inches(2), Inches(8), Inches(5)
                )

            tf = body_shape.text_frame
            tf.clear()
            for explanation in item.get("explanations", []):
                # Title
                p = tf.add_paragraph()
                p.text = (
                    f"{explanation.get('emoji', '')} {explanation.get('title', '')}"
                )
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = theme_color
                p.level = 0

                # Description
                p = tf.add_paragraph()
                p.text = explanation.get("description", "")
                p.font.size = Pt(24)
                p.level = 1

        elif item["type"] == "TESTIMONIALS":
            testimonials = item.get("testimonials", [])
            for testimonial in testimonials:
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)

                # Title
                slide.shapes.title.text = "Testimonial"
                slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = theme_color

                # Access the content placeholder
                body_shape = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                        body_shape = shape
                        break
                if body_shape is None:
                    body_shape = slide.shapes.add_textbox(
                        Inches(1), Inches(2), Inches(8), Inches(5)
                    )

                tf = body_shape.text_frame
                tf.clear()

                # Testimonial text
                p = tf.add_paragraph()
                p.text = f'"{testimonial.get("testimonial", "")}"'
                p.font.size = Pt(24)
                p.font.italic = True
                p.font.color.rgb = RGBColor(89, 89, 89)  # Dark grey
                p.level = 0

                # Name
                p = tf.add_paragraph()
                p.text = f"- {testimonial.get('firstName', '')} {testimonial.get('lastName', '')}"
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = theme_color
                p.level = 0

        elif item["type"] == "CTA":
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Add background color
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = theme_color

            # Add headline
            if item.get("headline"):
                left = Inches(1)
                top = Inches(2)
                width = prs.slide_width - Inches(2)
                height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = item["headline"]
                tf.paragraphs[0].font.size = Pt(44)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Add description
            if item.get("description"):
                left = Inches(1)
                top = Inches(3.5)
                width = prs.slide_width - Inches(2)
                height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = item["description"]
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Add Call to Action button
            if item.get("link"):
                # Create a rectangle shape as a button
                left = (prs.slide_width - Inches(3)) / 2
                top = Inches(5)
                width = Inches(3)
                height = Inches(0.8)
                button = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
                )
                button.fill.solid()
                button.fill.fore_color.rgb = RGBColor(255, 255, 255)
                button.line.color.rgb = RGBColor(255, 255, 255)

                # Add text to the button
                tf = button.text_frame
                p = tf.paragraphs[0]
                p.text = "Sign Up Now"
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = theme_color
                p.alignment = PP_ALIGN.CENTER

            # Add homepage link as text
            if item.get("homepageLink"):
                left = Inches(1)
                top = Inches(6)
                width = prs.slide_width - Inches(2)
                height = Inches(0.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = f"Visit our homepage: {item['homepageLink']}"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER

        # Add logo to each slide, ensure it's on top
        if logo_img:
            # Position logo at top-left corner
            try:
                logo = slide.shapes.add_picture(
                    logo_img, Inches(0.2), Inches(0.2), width=Inches(1.5)
                )
                # Bring logo to front
                logo.z_order = 0  # Ensure it's on top
            except:
                pass

    # Save the presentation to a file
    output_filename = f"decks/{uuid}.pptx"
    prs.save(output_filename)
    return output_filename
